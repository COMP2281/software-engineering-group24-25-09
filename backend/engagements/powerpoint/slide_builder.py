from pptx.shapes.autoshape import Shape
from pptx.slide import Slide as pptxSlide
from typing import cast
from engagements.powerpoint.shape_type import ShapeType


class SlideBuilder:
    @staticmethod
    def _find_shapes(template: pptxSlide) -> dict[ShapeType, int]:
        shapes = {}
        for shape in template.shapes:
            shape: Shape
            if shape.has_text_frame and shape.text in ShapeType:
                shape_type = cast(ShapeType, ShapeType[shape.text])
                shapes[shape_type] = shape.shape_id
        return shapes

    def __init__(self, template: pptxSlide) -> None:
        self.template = template
        self.shapes = self._find_shapes(template)

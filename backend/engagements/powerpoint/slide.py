from pptx import Presentation
from pptx.shapes.autoshape import Shape

file = open("templates.pptx", "rb")
slideshow = Presentation(file)
slides = slideshow.slides
for slide in slides:
    for shape in slide.shapes:
        shape: Shape
        print(shape.shape_type.name)
        if shape.shape_type.name == "PLACEHOLDER":
            print(shape.placeholder_format.type.name)
        print(shape.shape_id)
        print(shape.has_text_frame)
        if shape.has_text_frame:
            print(shape.text)
        print("\n")

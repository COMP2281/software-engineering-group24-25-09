import os
from pptx import Presentation
from pptx.slide import Slides
from engagements.slideshow.engagement_slide import EngagementSlide
from engagements.slideshow.placeholder_type import add_to_slide


class SlideshowBuilder:
    def reset(self) -> None:
        self.presentation = Presentation(self.layouts_file_path)

    def __init__(self, data_path: str) -> None:
        self.save_file_path = os.path.join(data_path, "export.pptx")
        self.layouts_file_path = os.path.join(data_path, "layouts.pptx")
        self.presentation = Presentation()
        self.reset()

    def save(self):
        self.presentation.save(self.save_file_path)

    def export(self, engagement_slides: list[EngagementSlide]) -> None:
        self.reset()
        for engagement_slide in engagement_slides:
            slides: Slides = self.presentation.slides
            slide = slides.add_slide(
                self.presentation.slide_layouts[engagement_slide.layout_index]
            )
            for placeholder_type in engagement_slide.values:
                placeholder_value = engagement_slide.values.get(placeholder_type)
                if placeholder_value:
                    add_to_slide[placeholder_type](slide, placeholder_value)
        self.save()

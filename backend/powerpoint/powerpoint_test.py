# CREATE SLIDE:
# Still not there yet
# Requires more robust ways of getting the actual content
# I don't know how images are being passed into the function I currently have it so it is taking in the root

# MERGE SLIDE:
# Pretty much done 
# Although the function does make all the other slides the same themes as the first one 
# which is true for powerpoint anyways as you can't have different themes on each slide
# This however could create a situation where one layout does not look well with another theme

# How to use CURRENTLY!
# create_slide(export_name(.pptx), title, text box, path to image<- might change)

# merge_presenation([list of slide names(.pptx)], export file name(.pptx))



from pptx import Presentation
import win32com.client
import os
import random # This is used for the slide theme(COULD LINK TO FRONTEND AND ALLOW USER TO CHOOSE THEME)
TEMPLATES = os.getenv("TEMPLATES", "templates") # Gets the path to the files
TEMPLATE_LIST = os.listdir(TEMPLATES) # Gets the files in the location

# From Existing temples create a slide
def create_slide(export_name, title_name, text, content):
    # Select a random template created
    selected_item = random.choice(TEMPLATE_LIST)

    # Open that slide
    os.chdir("templates")
    prs = Presentation(selected_item)
    slide = prs.slides[0]
    os.chdir("..")

    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))
        shape_format = shape.placeholder_format
        shape_type = str(shape_format.type).split(' (')[0]
        shape_index = shape.placeholder_format.idx
        placeholder = slide.placeholders[shape_index]
        print(shape_type)
        # Change title 
        if shape_type == "TITLE":
            placeholder.text = title_name

        # Change text
        elif shape_type == "BODY":
            placeholder.text = text     

        # Change content
        elif shape_type == "PICTURE":
            # Will probs need to change this in the final version but can stay for tests
            # I'm talking about using the os to get places rather than actual roots
            os.chdir("images")
            placeholder.insert_picture(content)
            os.chdir("..")

    
    # Save the changed slide
    os.chdir("created_slides")
    prs.save(export_name)
    os.chdir("..")


# Group the slides together
def merge_presentations(presentations, path):
  os.chdir("created_slides")
  ppt_instance = win32com.client.Dispatch('PowerPoint.Application')
  prs = ppt_instance.Presentations.open(os.path.abspath(presentations[0]), True, False, False)

  for i in range(1, len(presentations)):
      prs.Slides.InsertFromFile(os.path.abspath(presentations[i]), prs.Slides.Count)

  prs.SaveAs(os.path.abspath(path))
  prs.Close()


# PowerPoint creation
create_slide("test3.pptx", "Test", "Content would be here Why bullet points", "test.png")


merge_presentations(["test2.pptx", "test.pptx", "test3.pptx"],"FINAL.pptx")


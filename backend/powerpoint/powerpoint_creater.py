# How to use CURRENTLY!
# create_slide(export_name(.pptx) path, title, text box, path to image, template path)
# merge_presenation([list of slide names path(.pptx)], export file name path(.pptx))

from pptx import Presentation
import win32com.client
import os
# This is used for the slide theme(COULD LINK TO FRONTEND AND ALLOW USER TO CHOOSE THEME)
# import random 
# Only a idea for random template if none is chosen:
# TEMPLATES = os.getenv("TEMPLATES", "templates") # Gets the path to the files
# TEMPLATE_LIST = os.listdir(TEMPLATES) # Gets the files in the location

# From Existing temples create a slide
def create_slide(export_name, title_name, text, persons, content, template_choice):
    if os.path.isfile(export_name):
        print("File already exists") # LINK ERROR MESSAGE TO FRONT END HERE
    text_use = True # This is used to track if the main text is placed into the slide yet
    # A location to the slide choice is passed in, Could be made so that the file is passed in directly
    f = open(template_choice, "rb")
    prs = Presentation(f)
    f.close()
    # Open that slide
    slide = prs.slides[0]

    for shape in slide.placeholders:
        shape_format = shape.placeholder_format
        shape_type = str(shape_format.type).split(' (')[0]
        shape_index = shape.placeholder_format.idx
        placeholder = slide.placeholders[shape_index]
        # Change title 
        if shape_type == "TITLE":
            placeholder.text = title_name

        # Change text
        elif shape_type == "BODY":
            # Content area
            if text_use:
                placeholder.text = text
                text_use = False     
            # Person area
            elif len(persons) > 0:
                tf = placeholder.text_frame
                tf.text = "People involved:"
                for person in persons:
                    p = tf.add_paragraph()
                    p.text = "- " + person
                    p.level = 1
                
        # Change content
        elif shape_type == "PICTURE":
            # Get placeholder's original position and size
            left = placeholder.left
            top = placeholder.top
            width = placeholder.width
            height = placeholder.height

            # Remove placeholder before inserting image 
            sp = placeholder._element
            sp.getparent().remove(sp)

            # Insert image at the same position
            slide.shapes.add_picture(content, left, top, width=width, height=height)

    # Save the changed slide
    prs.save(export_name)

def edit_slide_text(slide_name, content, change):
    text_use = True
    # A location to the slide choice is passed in, Could be made so that the file is passed in directly
    f = open(slide_name, "rb")
    prs = Presentation(f)
    f.close()
    # Open that slide
    slide = prs.slides[0]

    for shape in slide.placeholders:
        shape_format = shape.placeholder_format
        shape_type = str(shape_format.type).split(' (')[0]
        shape_index = shape.placeholder_format.idx
        placeholder = slide.placeholders[shape_index]
    
        # Change title 
        if shape_type == "TITLE" and change == "title":
            placeholder.text = content
                # Change text
        elif shape_type == "BODY" and (change == "person" or change == "content"):
            # Content area
            if text_use and change == "content":
                placeholder.text = content
                text_use = False
            # Person area
            elif len(content) > 0 and change == "person" and not text_use:
                
                tf = placeholder.text_frame
                tf.text = "People involved:"
                for person in content:
                    p = tf.add_paragraph()
                    p.text = "- " + person
                    p.level = 1
            else:
                text_use = False
        
    # Save the changed slide
    prs.save(slide_name)

# Group the slides together
def merge_presentations(presentations, path):
  ppt_instance = win32com.client.Dispatch('PowerPoint.Application')
  prs = ppt_instance.Presentations.open(os.path.abspath(presentations[0]), True, False, False)

  for i in range(1, len(presentations)):
      prs.Slides.InsertFromFile(os.path.abspath(presentations[i]), prs.Slides.Count)

  prs.SaveAs(os.path.abspath(path))
  prs.Close()


# PowerPoint creation
# Tests
# create_slide(r"C:\Users\samda\Uni\Year 2\Software Engineering\software-engineering-group24-25-09\backend\powerpoint\created_slides\test.pptx",
#               "Test", "Content would be here Why bullet points", ["God", "John", "Jesus"], 
#               r"C:\Users\samda\Uni\Year 2\Software Engineering\software-engineering-group24-25-09\backend\powerpoint\images\test.png", 
#               r"C:\Users\samda\Uni\Year 2\Software Engineering\software-engineering-group24-25-09\backend\powerpoint\templates\template1Final.pptx")


# merge_presentations([r"C:\Users\samda\Uni\Year 2\Software Engineering\software-engineering-group24-25-09\backend\powerpoint\created_slides\test.pptx"],r"C:\Users\samda\Uni\Year 2\Software Engineering\software-engineering-group24-25-09\backend\powerpoint\created_slides\FINAL.pptx")

edit_slide_text(r"C:\Users\samda\Uni\Year 2\Software Engineering\software-engineering-group24-25-09\backend\powerpoint\created_slides\FINAL.pptx", r"C:\Users\samda\Uni\Year 2\Software Engineering\software-engineering-group24-25-09\backend\powerpoint\images\test2.jpg", "picture")
import os
from pptx import Presentation
from pptx.shapes.placeholder import LayoutPlaceholder
from pptx.slide import Slides, Slide
from engagements.powerpoint.engagement_slide import EngagementSlide


class SlideshowBuilder:
    def reset(self) -> None:
        self.presentation = Presentation(self.layouts_file_path)

    def __init__(self, data_path: str) -> None:
        self.save_file_path = os.path.join(data_path, "export.pptx")
        self.layouts_file_path = os.path.join(data_path, "layouts.pptx")
        self.presentation = Presentation()
        self.reset()

    def save(self):
        self.presentation.save(self.save_file_path)

    def export(self, engagement_slides: list[EngagementSlide]) -> None:
        self.reset()
        for engagement_slide in engagement_slides:
            slides: Slides = self.presentation.slides
            slide = slides.add_slide(
                self.presentation.slide_layouts[engagement_slide.layout_index]
            )
            for placeholder_type in engagement_slide.placeholders:
                placeholder_index = engagement_slide.placeholders[placeholder_type]
                placeholder_value = engagement_slide.values.get(placeholder_type)
                if placeholder_value:
                    slide_placeholder: LayoutPlaceholder = slide.placeholders[
                        placeholder_index
                    ]
                    slide_placeholder.text = str(placeholder_value)
        self.save()


slideshow_maker = SlideshowBuilder(".")
slides = [
    EngagementSlide()
    .set_title("hello john i hope youre happy")
    .set_summary(
        [
            "here is a summary of this slide",
            "john youre in for an incredible application",
        ]
    )
    .set_employees(["maks nowak", "god", "jesus"]),
    EngagementSlide()
    .set_title("slide 2 i hate this")
    .set_summary(
        "summarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummarysummary"
    )
    .set_employees(["ed jackson", "oscar ryley"]),
]
slideshow_maker.export(slides)

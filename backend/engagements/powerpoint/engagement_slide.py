from typing import cast, Self
from bs4 import BeautifulSoup
from pptx.shapes.placeholder import LayoutPlaceholder
from pptx.slide import SlideLayout
from engagements.powerpoint.types import PlaceholderType, PlaceholderDataType


class EngagementSlide:
    def _set_placeholders(self):
        self.placeholders = {}
        if not self.layout:
            raise Exception("Missing layout")
        for placeholder in self.layout.placeholders:
            placeholder: LayoutPlaceholder
            if placeholder.has_text_frame and placeholder.text in PlaceholderType:
                self.placeholders[
                    cast(PlaceholderType, PlaceholderType[placeholder.text])
                ] = placeholder.placeholder_format.idx

    def reset(self, layout: SlideLayout) -> Self:
        self.layout = layout
        self._set_placeholders()
        return self

    def __init__(self, layout: SlideLayout) -> None:
        self.layout: SlideLayout | None = None
        self.placeholders: dict[PlaceholderType, int] = {}
        self.placeholder_values: dict[PlaceholderType, PlaceholderDataType] = {}
        self.reset(layout)

    def get_layout(self) -> SlideLayout:
        return self.layout

    @property
    def title(self) -> str | None:
        return self.placeholder_values.get(PlaceholderType.TITLE)

    @title.setter
    def title(self, title: str) -> None:
        if PlaceholderType.TITLE in self.placeholders:
            self.placeholder_values[PlaceholderType.TITLE] = title

    @property
    def summary(self) -> list[str] | None:
        return self.placeholder_values.get(PlaceholderType.SUMMARY)

    @summary.setter
    def summary(self, summary: list[str]) -> None:
        if PlaceholderType.SUMMARY in self.placeholders:
            self.placeholder_values[PlaceholderType.SUMMARY] = summary

    @property
    def employees(self) -> list[str] | None:
        return self.placeholder_values.get(PlaceholderType.EMPLOYEES)

    @employees.setter
    def employees(self, employees: list[str]) -> None:
        if PlaceholderType.EMPLOYEES in self.placeholders:
            self.placeholder_values[PlaceholderType.EMPLOYEES] = employees

    @property
    def image(self) -> str | None:
        return self.placeholder_values.get(PlaceholderType.IMAGE)

    @image.setter
    def image(self, image: BeautifulSoup) -> None:
        # TODO: handle srcset attribute
        if PlaceholderType.IMAGE in self.placeholders:
            self.placeholder_values[PlaceholderType.IMAGE] = image["src"]

    def set_title(self, title: str) -> Self:
        self.title = title
        return self

    def set_summary(self, summary: list[str]) -> Self:
        self.summary = summary
        return self

    def set_employees(self, employees: list[str]) -> Self:
        self.employees = employees
        return self

    def set_image(self, image: BeautifulSoup) -> Self:
        self.image = image
        return self

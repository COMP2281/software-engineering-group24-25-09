from enum import Enum
from io import BytesIO
from typing import TypeAlias, Callable
from engagements.web import get
from pptx.shapes.placeholder import LayoutPlaceholder, PicturePlaceholder
from pptx.slide import Slide


class PlaceholderType(Enum):
    TITLE = "TITLE"
    SUMMARY = "SUMMARY"
    EMPLOYEES = "EMPLOYEES"
    IMAGE = "IMAGE"


PlaceholderDataType: TypeAlias = str | list[str]


def add_title_to_slide(slide: Slide, title: str) -> None:
    for placeholder in slide.slide_layout.placeholders:
        placeholder: LayoutPlaceholder
        if (
            placeholder.has_text_frame
            and placeholder.text == PlaceholderType.TITLE.value
        ):
            idx = placeholder.placeholder_format.idx
            slide.placeholders[idx].text = title


def add_summary_to_slide(slide: Slide, summary: list[str]) -> None:
    for placeholder in slide.slide_layout.placeholders:
        placeholder: LayoutPlaceholder
        if (
            placeholder.has_text_frame
            and placeholder.text == PlaceholderType.SUMMARY.value
        ):
            idx = placeholder.placeholder_format.idx
            slide.placeholders[idx].text = "\n".join(summary)


def add_employees_to_slide(slide: Slide, employees: list[str]) -> None:
    for placeholder in slide.slide_layout.placeholders:
        placeholder: LayoutPlaceholder
        if (
            placeholder.has_text_frame
            and placeholder.text == PlaceholderType.EMPLOYEES.value
        ):
            idx = placeholder.placeholder_format.idx
            slide.placeholders[idx].text = "\n".join(employees)


def add_image_to_slide(slide: Slide, image_url: str) -> None:
    for placeholder in slide.slide_layout.placeholders:
        placeholder: LayoutPlaceholder
        if (
            placeholder.has_text_frame
            and placeholder.text == PlaceholderType.IMAGE.value
        ):
            idx = placeholder.placeholder_format.idx
            response = get(image_url)
            image = BytesIO(response.content)
            slide.placeholders[idx].insert_picture(image)


add_to_slide: dict[PlaceholderType, Callable[[Slide, PlaceholderDataType], None]] = {
    PlaceholderType.TITLE: add_title_to_slide,
    PlaceholderType.SUMMARY: add_summary_to_slide,
    PlaceholderType.EMPLOYEES: add_employees_to_slide,
    PlaceholderType.IMAGE: add_image_to_slide,
}
